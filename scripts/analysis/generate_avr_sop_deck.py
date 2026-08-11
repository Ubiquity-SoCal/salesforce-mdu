"""
Generate the AVR Process SOP deck (PowerPoint).

AVR = Address Verification Request. Requestor emails an address -> Salesforce
auto-creates a Case in the Address Management app -> UBQ/GIS team verifies it
against Vetro and the address is ADDED, MODIFIED, or REJECTED -> a status-change
auto-email notifies the requestor -> everything is tracked in list views + a
Summary dashboard.

Design: SalesForce/docs/avr-sop-deck-design.md

Re-runnable: pulls the 7 screenshots from the Desktop AVRs folder, writes the
deck to SalesForce/data/output/, and drops a copy next to the screenshots.

Usage:  python SalesForce/scripts/analysis/generate_avr_sop_deck.py
"""
from __future__ import annotations

import os
import shutil
import sys

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

# --------------------------------------------------------------------------- #
# Paths
# --------------------------------------------------------------------------- #
HOME = os.path.expanduser("~")
SHOTS_DIR = os.path.join(HOME, "OneDrive - Ubiquity Management", "Desktop", "AVRs")
REPO_OUT = os.path.join(
    HOME, "Work_Projects", "SalesForce", "data", "output", "avr-process-sop.pptx"
)

IMG = {
    "email": "AVR Email received (Creates Case Record).png",
    "case": "AVR Case Record.png",
    "status": "AVR Status.png",
    "auto_email": "AVR Auto Email when Status Changes on Case.png",
    "add_confirm": "AVR - Address Add Confirmation.png",
    "open_cases": "AVR Default Pipelineview -Open Cases.png",
    "all_cases": "AVR Pipelineview -All Cases.png",
    "summary": "AVR Summary.png",
}


def shot(key: str) -> str:
    p = os.path.join(SHOTS_DIR, IMG[key])
    if not os.path.isfile(p):
        sys.exit(f"ERROR: screenshot not found: {p}")
    return p


# --------------------------------------------------------------------------- #
# Palette & type
# --------------------------------------------------------------------------- #
NAVY = RGBColor(0x03, 0x2D, 0x60)     # Salesforce navy (titles)
BLUE = RGBColor(0x01, 0x76, 0xD3)     # Salesforce brand blue (accent)
LIGHT_BLUE = RGBColor(0xEA, 0xF5, 0xFE)
GRAY = RGBColor(0x6B, 0x72, 0x80)     # muted support text
DARK = RGBColor(0x18, 0x1D, 0x27)     # body text
HAIRLINE = RGBColor(0xD5, 0xDB, 0xE1)  # image border / dividers
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

GREEN = RGBColor(0x2E, 0x84, 0x4A)    # Added
RED = RGBColor(0xEA, 0x00, 0x1E)      # Rejected

FONT = "Segoe UI"

# Slide geometry (16:9)
SW, SH = Inches(13.333), Inches(7.5)
MARGIN = Inches(0.6)
CONTENT_W = SW - 2 * MARGIN

FOOTER_TXT = "AVR Process · Address Verification Requests"


# --------------------------------------------------------------------------- #
# Low-level helpers
# --------------------------------------------------------------------------- #
def _set_text(tf, runs, align=PP_ALIGN.LEFT, space_after=None, line_spacing=None):
    """runs: list of (text, size, color, bold) tuples for a single paragraph,
    or list-of-lists for multiple paragraphs."""
    if runs and not isinstance(runs[0], list):
        runs = [runs]
    tf.word_wrap = True
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if space_after is not None:
            p.space_after = space_after
        if line_spacing is not None:
            p.line_spacing = line_spacing
        for (text, size, color, bold) in para:
            r = p.add_run()
            r.text = text
            r.font.name = FONT
            r.font.size = Pt(size)
            r.font.color.rgb = color
            r.font.bold = bold


def textbox(slide, left, top, width, height, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    return tb, tf


def rect(slide, left, top, width, height, fill=None, line=None, line_w=None,
         shape=MSO_SHAPE.RECTANGLE, radius=None):
    sp = slide.shapes.add_shape(shape, left, top, width, height)
    sp.shadow.inherit = False
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = line_w or Pt(1)
    if radius is not None and shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        try:
            sp.adjustments[0] = radius
        except Exception:
            pass
    return sp


def fit_image(slide, path, box_l, box_t, box_w, box_h, border=True):
    """Place an image scaled to fit box (aspect preserved), centered, with a
    subtle border. Returns the Picture."""
    with Image.open(path) as im:
        iw, ih = im.size
    box_ratio = box_w / box_h
    img_ratio = iw / ih
    if img_ratio > box_ratio:          # image wider -> constrain width
        w = box_w
        h = Emu(int(box_w * ih / iw))
    else:                              # image taller -> constrain height
        h = box_h
        w = Emu(int(box_h * iw / ih))
    left = box_l + (box_w - w) // 2
    top = box_t + (box_h - h) // 2
    pic = slide.shapes.add_picture(path, left, top, w, h)
    if border:
        pic.line.color.rgb = HAIRLINE
        pic.line.width = Pt(0.75)
    return pic


# --------------------------------------------------------------------------- #
# Slide chrome
# --------------------------------------------------------------------------- #
def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def content_header(slide, kicker, title, page):
    """Left accent rule + kicker + title + footer. Returns bottom-y of header."""
    # left accent rule
    rect(slide, MARGIN, Inches(0.55), Inches(0.09), Inches(0.95), fill=BLUE)
    # kicker
    _, tf = textbox(slide, MARGIN + Inches(0.28), Inches(0.52),
                    CONTENT_W - Inches(0.28), Inches(0.3))
    _set_text(tf, [(kicker.upper(), 12, BLUE, True)])
    tf.paragraphs[0].runs[0].font.spacing = Pt(1)
    # title
    _, tf = textbox(slide, MARGIN + Inches(0.28), Inches(0.82),
                    CONTENT_W - Inches(0.28), Inches(0.7))
    _set_text(tf, [(title, 26, NAVY, True)])
    footer(slide, page)
    return Inches(1.75)


def footer(slide, page):
    # divider
    rect(slide, MARGIN, Inches(7.02), CONTENT_W, Emu(9525), fill=HAIRLINE)
    _, tf = textbox(slide, MARGIN, Inches(7.08), Inches(8), Inches(0.3))
    _set_text(tf, [(FOOTER_TXT, 9, GRAY, False)])
    _, tf = textbox(slide, SW - MARGIN - Inches(2), Inches(7.08),
                    Inches(2), Inches(0.3))
    _set_text(tf, [(f"{page}", 9, GRAY, False)], align=PP_ALIGN.RIGHT)


def outcome_chip(slide, left, top, width, height, color, label, sub):
    rect(slide, left, top, width, height, fill=WHITE, line=HAIRLINE,
         line_w=Pt(1), shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.08)
    rect(slide, left, top, Inches(0.09), height, fill=color)  # color spine
    _, tf = textbox(slide, left + Inches(0.28), top + Inches(0.12),
                    width - Inches(0.4), height - Inches(0.24),
                    anchor=MSO_ANCHOR.MIDDLE)
    _set_text(
        tf,
        [
            [(label, 14, color, True)],
            [(sub, 10.5, GRAY, False)],
        ],
        line_spacing=1.05,
    )


# --------------------------------------------------------------------------- #
# Slides
# --------------------------------------------------------------------------- #
def slide_title(prs):
    s = blank_slide(prs)
    # left navy band
    band_w = Inches(4.6)
    rect(s, 0, 0, band_w, SH, fill=NAVY)
    rect(s, band_w, 0, Inches(0.10), SH, fill=BLUE)
    # band content
    _, tf = textbox(s, Inches(0.7), Inches(2.55), band_w - Inches(1.1),
                    Inches(0.4))
    _set_text(tf, [("STANDARD OPERATING PROCEDURE", 12, RGBColor(0x9E, 0xC5, 0xF0), True)])
    _, tf = textbox(s, Inches(0.7), Inches(3.0), band_w - Inches(1.1), Inches(2.2))
    _set_text(
        tf,
        [
            [("AVR", 60, WHITE, True)],
            [("Process", 60, WHITE, True)],
        ],
        line_spacing=0.95,
    )
    # right content
    rx = band_w + Inches(0.9)
    _, tf = textbox(s, rx, Inches(2.75), SW - rx - MARGIN, Inches(1.0))
    _set_text(
        tf,
        [
            [("Address Verification Requests", 26, NAVY, True)],
            [("Address Management app · Salesforce", 15, GRAY, False)],
        ],
        line_spacing=1.15,
    )
    rect(s, rx, Inches(3.95), Inches(1.4), Pt(3), fill=BLUE)
    _, tf = textbox(s, rx, Inches(4.25), SW - rx - MARGIN, Inches(1.6))
    _set_text(
        tf,
        [
            [("How an address request flows from email to resolution:", 13, DARK, False)],
            [("added, modified, or rejected, and how it's tracked.", 13, DARK, False)],
        ],
        line_spacing=1.25,
    )
    _, tf = textbox(s, rx, Inches(6.6), SW - rx - MARGIN, Inches(0.4))
    _set_text(tf, [("As of July 1, 2026", 11, GRAY, False)])
    return s


def slide_what(prs):
    s = blank_slide(prs)
    content_header(s, "Overview", "What is an AVR?", 2)
    _, tf = textbox(s, MARGIN + Inches(0.28), Inches(1.75),
                    CONTENT_W - Inches(0.28), Inches(1.5))
    _set_text(
        tf,
        [
            [("An ", 15, DARK, False),
             ("Address Verification Request", 15, NAVY, True),
             (" is a request to confirm an address against Ubiquity's", 15, DARK, False)],
            [("network. A requestor (e.g. FiberFirst) submits an address; the ", 15, DARK, False),
             ("UBQ / GIS team", 15, NAVY, True),
             (" verifies it", 15, DARK, False)],
            [("against Vetro and resolves it one of three ways:", 15, DARK, False)],
        ],
        line_spacing=1.25,
    )
    # three outcome chips
    top = Inches(3.35)
    h = Inches(1.5)
    gap = Inches(0.35)
    w = (CONTENT_W - 2 * gap) / 3
    outcome_chip(s, MARGIN, top, w, h, GREEN, "Added",
                 "New address added\nto Vetro / the network")
    outcome_chip(s, MARGIN + w + gap, top, w, h, BLUE, "Modified",
                 "Existing address\ncorrected or reformatted")
    outcome_chip(s, MARGIN + 2 * (w + gap), top, w, h, RED, "Rejected",
                 "Not serviceable, invalid,\nor not in network")
    # footnote
    _, tf = textbox(s, MARGIN, Inches(5.4), CONTENT_W, Inches(1.0))
    _set_text(
        tf,
        [
            [("Where it lives:  ", 12, NAVY, True),
             ("Tracked as Cases in the Address Management app in Salesforce. "
              "Requests arrive by email and", 12, DARK, False)],
            [("create a Case automatically; a status change emails the requestor back. "
              "Added + Modified count as", 12, DARK, False)],
            [("“Complete” on the dashboard; Rejected counts as “Invalid.”", 12, DARK, False)],
        ],
        line_spacing=1.2,
    )
    return s


def _flow_box(slide, left, top, width, height, n, label):
    rect(slide, left, top, width, height, fill=WHITE, line=BLUE, line_w=Pt(1.25),
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.12)
    # number badge
    b = Inches(0.34)
    rect(slide, left + Inches(0.12), top + Inches(0.12), b, b, fill=BLUE,
         shape=MSO_SHAPE.OVAL)
    _, tf = textbox(slide, left + Inches(0.12), top + Inches(0.12), b, b,
                    anchor=MSO_ANCHOR.MIDDLE)
    _set_text(tf, [(str(n), 12, WHITE, True)], align=PP_ALIGN.CENTER)
    # label
    _, tf = textbox(slide, left + Inches(0.1), top + Inches(0.5),
                    width - Inches(0.2), height - Inches(0.6),
                    anchor=MSO_ANCHOR.MIDDLE)
    _set_text(tf, [(label, 11.5, NAVY, True)], align=PP_ALIGN.CENTER,
              line_spacing=1.0)


def _arrow(slide, left, top, width, height):
    a = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left, top, width, height)
    a.shadow.inherit = False
    a.fill.solid()
    a.fill.fore_color.rgb = RGBColor(0xB6, 0xD8, 0xF5)
    a.line.fill.background()
    try:
        a.adjustments[0] = 0.55
        a.adjustments[1] = 0.55
    except Exception:
        pass


def slide_lifecycle(prs):
    s = blank_slide(prs)
    content_header(s, "The AVR Lifecycle", "From email in to resolution tracked", 3)
    stages = [
        "Request in\n(email)",
        "Case created\nautomatically",
        "UBQ / GIS\nreview",
        "Outcome set\n(status)",
        "Requestor\nauto-emailed",
        "Added to\nBE + FSM",
        "Tracked on\ndashboard",
    ]
    n = len(stages)
    arrow_w = Inches(0.34)
    box_w = (CONTENT_W - (n - 1) * arrow_w) / n
    box_h = Inches(1.25)
    top = Inches(2.15)
    x = MARGIN
    for i, label in enumerate(stages):
        _flow_box(s, x, top, box_w, box_h, i + 1, label)
        x += box_w
        if i < n - 1:
            _arrow(s, x, top + box_h / 2 - Inches(0.18), arrow_w, Inches(0.36))
            x += arrow_w

    # branch from stage 4 (Outcome) into three chips
    # stage 4 center x
    s4_left = MARGIN + 3 * (box_w + arrow_w)
    s4_cx = s4_left + box_w / 2
    chips_top = Inches(4.35)
    # connector line down
    rect(s, s4_cx - Emu(6350), top + box_h, Emu(12700),
         chips_top - (top + box_h), fill=RGBColor(0xB6, 0xD8, 0xF5))
    _, tf = textbox(s, MARGIN, Inches(3.85), CONTENT_W, Inches(0.35))
    _set_text(tf, [("The outcome is always one of three:", 13, GRAY, True)],
              align=PP_ALIGN.CENTER)
    ch = Inches(1.35)
    gap = Inches(0.35)
    cw = (CONTENT_W - 2 * gap) / 3
    outcome_chip(s, MARGIN, chips_top, cw, ch, GREEN, "Added",
                 "Uploaded to Vetro /\nPending Upload to Vetro")
    outcome_chip(s, MARGIN + cw + gap, chips_top, cw, ch, BLUE, "Modified",
                 "Format Change in Vetro /\nOmnia/CHR Changes Required")
    outcome_chip(s, MARGIN + 2 * (cw + gap), chips_top, cw, ch, RED, "Rejected",
                 "Invalid / Not in Network /\nUnserviceable")
    return s


def image_slide(prs, kicker, title, page, img_key, bullets,
                layout="top", img_frac=0.58, top_img_w=None):
    """Generic step slide. layout 'top' = bullets then wide image below;
    'side' = image right, bullets left. top_img_w caps/centers the image
    width for the 'top' layout (defaults to full content width)."""
    s = blank_slide(prs)
    content_header(s, kicker, title, page)
    body_top = Inches(1.75)
    body_bottom = Inches(6.9)
    if layout == "top":
        # bullets
        bh = Inches(1.55)
        _bullets(s, MARGIN, body_top, CONTENT_W, bh, bullets)
        img_t = body_top + bh + Inches(0.1)
        box_w = top_img_w or CONTENT_W
        box_l = MARGIN + (CONTENT_W - box_w) // 2
        fit_image(s, shot(img_key), box_l, img_t, box_w,
                  body_bottom - img_t)
    else:  # side
        img_w = CONTENT_W * img_frac
        txt_w = CONTENT_W - img_w - Inches(0.4)
        _bullets(s, MARGIN, body_top, txt_w, body_bottom - body_top, bullets)
        fit_image(s, shot(img_key), MARGIN + txt_w + Inches(0.4), body_top,
                  img_w, body_bottom - body_top)
    return s


def _bullets(slide, left, top, width, height, bullets):
    _, tf = textbox(slide, left, top, width, height)
    paras = []
    for b in bullets:
        if isinstance(b, tuple):  # (lead, rest)
            paras.append([("•  ", 13, BLUE, True),
                          (b[0], 13, NAVY, True),
                          (b[1], 13, DARK, False)])
        else:
            paras.append([("•  ", 13, BLUE, True), (b, 13, DARK, False)])
    _set_text(tf, paras, space_after=Pt(9), line_spacing=1.12)


def slide_status(prs):
    s = blank_slide(prs)
    content_header(s, "Step 3 · Review & Outcome",
                   "UBQ / GIS reviews and sets the AVR Status", 6)
    body_top = Inches(1.8)
    # image left (small picklist)
    img_w = Inches(3.9)
    fit_image(s, shot("status"), MARGIN, body_top, img_w, Inches(4.9))
    _, tf = textbox(s, MARGIN, Inches(1.82), img_w, Inches(0.3))
    _set_text(tf, [("The AVR Status picklist", 11, GRAY, True)],
              align=PP_ALIGN.CENTER)

    # right: three grouped outcome cards
    gx = MARGIN + img_w + Inches(0.5)
    gw = SW - MARGIN - gx
    groups = [
        (GREEN, "Added", ["Uploaded to Vetro", "Pending Upload to Vetro"]),
        (BLUE, "Modified", ["Format Change Completed in Vetro",
                            "Omnia/CHR Changes Required"]),
        (RED, "Rejected", ["Invalid Address Request", "Address Not in Network",
                           "Future / Unserviceable",
                           "See “Engineering / Add. Management Notes”"]),
    ]
    y = body_top
    for color, name, items in groups:
        card_h = Inches(0.44) + Inches(0.30) * len(items) + Inches(0.18)
        rect(s, gx, y, gw, card_h, fill=WHITE, line=HAIRLINE, line_w=Pt(1),
             shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.06)
        rect(s, gx, y, Inches(0.09), card_h, fill=color)
        _, tf = textbox(s, gx + Inches(0.28), y + Inches(0.1),
                        gw - Inches(0.5), Inches(0.34))
        _set_text(tf, [(name.upper(), 12.5, color, True)])
        _, tf = textbox(s, gx + Inches(0.28), y + Inches(0.46),
                        gw - Inches(0.5), card_h - Inches(0.5))
        _set_text(
            tf,
            [[("•  ", 11.5, color, True), (it, 11.5, DARK, False)] for it in items],
            line_spacing=1.05, space_after=Pt(2),
        )
        y = Emu(int(y) + int(card_h) + int(Inches(0.16)))
    # mid-workflow note
    _, tf = textbox(s, gx, y + Inches(0.02), gw, Inches(0.6))
    _set_text(
        tf,
        [[("Mid-workflow (not a final outcome):  ", 10.5, NAVY, True),
          ("Pending Review Board · Utilize Master List · Further Investigation Required",
           10.5, GRAY, False)]],
        line_spacing=1.1,
    )
    return s


def slide_pipeline(prs):
    s = blank_slide(prs)
    content_header(s, "Tracking · Pipeline", "Cases are tracked in list views", 9)
    _bullets(
        s, MARGIN, Inches(1.75), CONTENT_W, Inches(1.15),
        [
            ("Default view: ", "“All Open Cases” (pinned). Other list views: All Cases, "
             "Invalid Requests, My Cases, My Open Cases, Recently Viewed."),
            ("All Cases: ", "every AVR with Case #, address type, AVR #, requested address, "
             "state, status, requestor, opened date, engineering notes & HQ reviewed date."),
        ],
    )
    # two wide images stacked
    fit_image(s, shot("open_cases"), MARGIN, Inches(3.0), CONTENT_W, Inches(1.85))
    fit_image(s, shot("all_cases"), MARGIN, Inches(5.0), CONTENT_W, Inches(1.7))
    return s


def slide_summary(prs):
    s = blank_slide(prs)
    content_header(s, "Tracking · Dashboard", "AVR Summary dashboard", 10)
    # stat chips
    stats = [
        ("224", "AVRs tracked"),
        ("2.37 days", "avg turnaround"),
        ("SFU 2.3 · MDU 4.5 · BUS 1.6", "avg days by type"),
        ("Complete vs Invalid", "simplified buckets"),
    ]
    n = len(stats)
    gap = Inches(0.3)
    cw = (CONTENT_W - (n - 1) * gap) / n
    top = Inches(1.7)
    ch = Inches(0.95)
    x = MARGIN
    for val, lab in stats:
        rect(s, x, top, cw, ch, fill=LIGHT_BLUE, line=None,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.1)
        _, tf = textbox(s, x + Inches(0.15), top + Inches(0.12), cw - Inches(0.3),
                        ch - Inches(0.24), anchor=MSO_ANCHOR.MIDDLE)
        _set_text(
            tf,
            [[(val, 17, NAVY, True)], [(lab, 10.5, GRAY, False)]],
            align=PP_ALIGN.CENTER, line_spacing=1.0,
        )
        x += cw + gap
    # dashboard image
    fit_image(s, shot("summary"), MARGIN, Inches(2.85), CONTENT_W, Inches(4.05))
    return s


# --------------------------------------------------------------------------- #
# Build
# --------------------------------------------------------------------------- #
def build():
    prs = Presentation()
    prs.slide_width = SW
    prs.slide_height = SH

    slide_title(prs)
    slide_what(prs)
    slide_lifecycle(prs)
    image_slide(
        prs, "Step 1 · Request In", "A request arrives by email", 4, "email",
        [
            ("The requestor emails ", "the address to verify (requestor, contact, full "
             "parsed address, location type SFU/MDU/BUS, region, lat/long, priority)."),
            ("Salesforce ingests the email ", "and auto-creates a Case (Case Origin = Email)."),
            ("The full email is retained ", "on the Case under “Email Details.”"),
        ],
        layout="top",
    )
    image_slide(
        prs, "Step 2 · Case Created", "The AVR Case record", 5, "case",
        [
            ("Each request becomes a Case ", "with its own AVR Number (e.g. AVR-92) "
             "and Case Number."),
            ("Case Information / Request Description ", "capture who asked, the issue, "
             "status, and key dates."),
            ("Address Details ", "hold the parsed address + lat/long; ",),
            ("UBQ Review ", "is where the team records the outcome, engineering notes, "
             "reviewer, and # of addresses added to Vetro."),
        ],
        layout="top",
    )
    slide_status(prs)
    image_slide(
        prs, "Step 4 · Requestor Notified",
        "A status change auto-emails the requestor", 7, "auto_email",
        [
            ("Automatic. ", "When the AVR Status changes on the Case, a Salesforce flow "
             "automatically emails the requestor."),
            ("The email includes ", "the AVR #, requested vs. updated address, the new "
             "status, HQ date reviewed, reviewer, and engineering / management notes."),
            ("Sent from the GIS Team ", "with the original requestor on To and stakeholders "
             "on Cc."),
        ],
        layout="side", img_frac=0.42,
    )
    image_slide(
        prs, "Step 5 · Added Downstream",
        "The verified address is added to COS Business Engine + FSM", 8,
        "add_confirm",
        [
            ("Once verified, ", "the address is added into the downstream systems: "
             "COS Business Engine (BE) and Field Service Management (FSM)."),
            ("Confirmed on-thread. ", "The requestor replies on the AVR thread confirming "
             "the add, so the address is live for billing and field service."),
            ("Loop closed. ", "Case, status notification, and downstream confirmation all "
             "stay together under the same AVR number."),
        ],
        layout="top", top_img_w=Inches(9.5),
    )
    slide_pipeline(prs)
    slide_summary(prs)

    os.makedirs(os.path.dirname(REPO_OUT), exist_ok=True)
    prs.save(REPO_OUT)
    print(f"Saved: {REPO_OUT}  ({len(prs.slides)} slides)")

    # drop a copy next to the screenshots
    desktop_copy = os.path.join(SHOTS_DIR, "AVR Process SOP.pptx")
    shutil.copyfile(REPO_OUT, desktop_copy)
    print(f"Copied: {desktop_copy}")


if __name__ == "__main__":
    build()
